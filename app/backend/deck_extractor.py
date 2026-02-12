import re
from dataclasses import dataclass
from pathlib import Path
from typing import List

from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_DECK_EXTENSIONS = {".pdf", ".pptx", ".ppt"}


@dataclass
class DeckExtractionResult:
    extracted_text: str
    extracted_json: List[dict]
    num_pages_or_slides: int


def detect_extension(filename: str) -> str:
    return Path(filename or "").suffix.lower()


def sanitize_filename(filename: str) -> str:
    candidate = Path(filename or "").name
    if candidate in {"", ".", ".."}:
        candidate = "deck"

    sanitized = re.sub(r"[^A-Za-z0-9._-]", "_", candidate)
    if sanitized in {"", ".", ".."}:
        sanitized = "deck"

    stem = Path(sanitized).stem[:120] or "deck"
    ext = Path(sanitized).suffix[:20]
    return f"{stem}{ext}"


def validate_deck_extension(extension: str) -> None:
    if extension not in SUPPORTED_DECK_EXTENSIONS:
        raise ValueError("Unsupported deck format. Please upload PDF, PPTX, or PPT.")
    if extension == ".ppt":
        raise ValueError("Legacy .ppt is not supported yet. Please upload PDF or PPTX.")


def extract_deck_text(deck_path: Path) -> DeckExtractionResult:
    extension = deck_path.suffix.lower()
    if extension == ".pdf":
        return _extract_pdf(deck_path)
    if extension == ".pptx":
        return _extract_pptx(deck_path)
    if extension == ".ppt":
        raise ValueError("Legacy .ppt is not supported yet. Please upload PDF or PPTX.")
    raise ValueError("Unsupported deck format. Please upload PDF or PPTX.")


def _extract_pdf(deck_path: Path) -> DeckExtractionResult:
    reader = PdfReader(str(deck_path))
    entries: List[dict] = []
    merged: List[str] = []

    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        entries.append({"index": index, "text": text})
        merged.append(f"PAGE {index}: {text}")

    return DeckExtractionResult(
        extracted_text="\n\n".join(merged).strip(),
        extracted_json=entries,
        num_pages_or_slides=len(reader.pages),
    )


def _extract_pptx(deck_path: Path) -> DeckExtractionResult:
    presentation = Presentation(str(deck_path))
    entries: List[dict] = []
    merged: List[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        text_chunks: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                text_chunks.append(text.strip())

        slide_text = "\n".join(chunk for chunk in text_chunks if chunk).strip()
        entries.append({"index": index, "text": slide_text})
        merged.append(f"SLIDE {index}: {slide_text}")

    return DeckExtractionResult(
        extracted_text="\n\n".join(merged).strip(),
        extracted_json=entries,
        num_pages_or_slides=len(presentation.slides),
    )
